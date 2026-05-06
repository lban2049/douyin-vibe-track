#!/usr/bin/env python3
"""Run one daily multi-platform account monitoring pass."""

from __future__ import annotations

import argparse
import concurrent.futures
import json
import re
import shutil
import subprocess
import sys
import time
from dataclasses import asdict, dataclass
from datetime import datetime, timedelta
from pathlib import Path
from typing import Any
from zoneinfo import ZoneInfo

import requests

from platform_core import (
    fetch_posts_page,
    get_session,
    load_json,
    normalize_platform,
    platform_label,
    should_retry_exception,
    summarize_counts,
    unwrap_posts_page,
    write_json,
)


SCRIPT_DIR = Path(__file__).resolve().parent
WORKSPACE_DIRNAME = "douyin-vibe-track"
DEFAULT_THRESHOLDS = {
    "douyin": {"like_count": 10000},
    "wechat_channels": {"like_count": 10000},
    "kuaishou": {"like_count": 10000},
}


@dataclass
class AccountFetchResult:
    account: dict[str, Any]
    candidates: list[Candidate]
    errors: list[RunIssue]
    duration_seconds: float
    posts_seen: int
    pages_fetched: int


@dataclass
class DownloadResult:
    item: Candidate
    error: RunIssue | None
    duration_seconds: float


@dataclass
class RuntimeMetrics:
    fetch_stage_seconds: float = 0.0
    download_stage_seconds: float = 0.0
    ppt_stage_seconds: float = 0.0
    total_seconds: float = 0.0
    fetched_accounts: int = 0
    failed_accounts: int = 0
    download_candidates: int = 0
    downloaded_videos: int = 0
    account_timings: list[dict[str, Any]] | None = None
    slow_accounts: list[dict[str, Any]] | None = None

    def __post_init__(self) -> None:
        if self.account_timings is None:
            self.account_timings = []
        if self.slow_accounts is None:
            self.slow_accounts = []


@dataclass
class Candidate:
    platform: str
    post_id: str
    author_name: str
    author_id: str
    desc: str
    share_url: str
    create_time: int
    metrics: dict[str, int]
    video: dict[str, Any]
    video_path: str = ""

    @property
    def like_count(self) -> int:
        return int((self.metrics or {}).get("like_count") or 0)

    @property
    def collect_count(self) -> int:
        return int((self.metrics or {}).get("collect_count") or 0)

    @property
    def share_count(self) -> int:
        return int((self.metrics or {}).get("share_count") or 0)

    @property
    def comment_count(self) -> int:
        return int((self.metrics or {}).get("comment_count") or 0)

    @property
    def video_urls(self) -> list[str]:
        return [str(url) for url in (self.video or {}).get("download_urls") or []]

    @property
    def needs_decrypt(self) -> bool:
        return bool((self.video or {}).get("needs_decrypt"))

    @property
    def decrypt_key(self) -> str:
        return str((self.video or {}).get("decrypt_key") or "")


@dataclass
class RunIssue:
    phase: str
    message: str
    fix: str
    account: str = ""
    platform: str = ""
    post_id: str = ""
    command: str = ""

    def to_dict(self) -> dict[str, str]:
        return {
            key: value
            for key, value in {
                "phase": self.phase,
                "platform": self.platform,
                "account": self.account,
                "post_id": self.post_id,
                "message": self.message,
                "fix": self.fix,
                "command": self.command,
            }.items()
            if value
        }


def shell_quote(value: str) -> str:
    return "'" + value.replace("'", "'\"'\"'") + "'"


def resolve_reports_root(workspace: Path, config: dict[str, Any]) -> Path:
    raw_reports_dir = str(config.get("reports_dir") or "").strip()
    if not raw_reports_dir:
        return workspace / "reports"
    reports_root = Path(raw_reports_dir).expanduser()
    if not reports_root.is_absolute():
        reports_root = workspace / reports_root
    return reports_root


def migrate_config(config: dict[str, Any]) -> tuple[dict[str, Any], bool]:
    updated = False
    config = dict(config)
    thresholds = config.get("thresholds")
    legacy_like_threshold = int(config.get("like_threshold") or DEFAULT_THRESHOLDS["douyin"]["like_count"])
    if not isinstance(thresholds, dict):
        thresholds = {}
        updated = True
    for platform in ("douyin", "wechat_channels", "kuaishou"):
        value = thresholds.get(platform)
        if not isinstance(value, dict):
            thresholds[platform] = {"like_count": legacy_like_threshold}
            updated = True
        elif "like_count" not in value:
            value["like_count"] = legacy_like_threshold
            updated = True
    config["thresholds"] = thresholds

    platforms = config.get("platforms")
    if not isinstance(platforms, dict):
        platforms = {}
        updated = True
    for platform in ("wechat_channels", "kuaishou"):
        platform_block = platforms.get(platform)
        if not isinstance(platform_block, dict):
            platforms[platform] = {"enabled": True}
            updated = True
        elif "enabled" not in platform_block:
            platform_block["enabled"] = True
            updated = True
    config["platforms"] = platforms

    wechat_config = config.get("wechat_channels")
    if not isinstance(wechat_config, dict):
        wechat_config = {}
        updated = True
    if "decrypt_mode" not in wechat_config:
        wechat_config["decrypt_mode"] = "builtin"
        updated = True
    if "keep_encrypted_copy" not in wechat_config:
        wechat_config["keep_encrypted_copy"] = False
        updated = True
    config["wechat_channels"] = wechat_config

    defaults = {
        "fetch_concurrency": 6,
        "download_concurrency": 3,
        "connect_timeout_seconds": 5,
        "read_timeout_seconds": int(config.get("request_timeout_seconds") or 30),
        "max_pages_per_account": 3,
        "slow_account_threshold_seconds": 45,
        "retryable_status_codes": [429, 500, 502, 503, 504],
    }
    for key, value in defaults.items():
        if key not in config:
            config[key] = value
            updated = True
    return config, updated


def migrate_account(account: dict[str, Any]) -> tuple[dict[str, Any], bool]:
    updated = False
    platform = normalize_platform(str(account.get("platform") or "douyin"))
    migrated = dict(account)
    migrated["platform"] = platform

    identity = migrated.get("identity")
    if not isinstance(identity, dict):
        identity = {}
        updated = True

    if platform in {"douyin", "tiktok"}:
        mappings = {
            "sec_user_id": migrated.get("sec_user_id") or "",
            "user_id": migrated.get("user_id") or "",
            "unique_id": migrated.get("unique_id") or "",
        }
    elif platform == "kuaishou":
        mappings = {
            "eid": migrated.get("eid") or migrated.get("user_id") or "",
            "user_id": migrated.get("numeric_user_id") or "",
            "unique_id": migrated.get("unique_id") or "",
        }
    else:
        mappings = {
            "username": migrated.get("username") or migrated.get("sec_user_id") or "",
        }
    for key, value in mappings.items():
        if value and not identity.get(key):
            identity[key] = str(value)
            updated = True
    migrated["identity"] = identity

    snapshot = migrated.get("profile_snapshot")
    if not isinstance(snapshot, dict):
        snapshot = {}
        updated = True
    snapshot_defaults = {
        "display_name": migrated.get("display_name") or "",
        "signature": migrated.get("signature") or "",
        "follower_count": summarize_counts(migrated.get("follower_count")),
        "following_count": summarize_counts(migrated.get("following_count")),
        "post_count": summarize_counts(migrated.get("aweme_count") or migrated.get("post_count")),
    }
    for key, value in snapshot_defaults.items():
        if key not in snapshot and value not in ("", 0):
            snapshot[key] = value
            updated = True
    migrated["profile_snapshot"] = snapshot
    migrated.setdefault("enabled", True)
    return migrated, updated


def migrate_accounts(accounts_data: dict[str, Any]) -> tuple[dict[str, Any], bool]:
    updated = False
    accounts = []
    for account in accounts_data.get("accounts", []):
        migrated, changed = migrate_account(account)
        accounts.append(migrated)
        updated = updated or changed
    return {"accounts": accounts}, updated


def migrate_state(state: dict[str, Any]) -> tuple[dict[str, Any], bool]:
    updated = False
    state = dict(state)
    post_keys = state.get("reported_post_keys")
    if not isinstance(post_keys, list):
        post_keys = []
        updated = True
    if state.get("reported_aweme_ids"):
        for aweme_id in state.get("reported_aweme_ids") or []:
            key = f"douyin:{aweme_id}"
            if key not in post_keys:
                post_keys.append(key)
                updated = True
    state["reported_post_keys"] = sorted({str(value) for value in post_keys if value})
    return state, updated


def account_label(account: dict[str, Any]) -> str:
    snapshot = account.get("profile_snapshot") or {}
    identity = account.get("identity") or {}
    return str(
        account.get("display_name")
        or snapshot.get("display_name")
        or identity.get("username")
        or identity.get("sec_user_id")
        or identity.get("eid")
        or identity.get("unique_id")
        or identity.get("user_id")
        or "unknown-account"
    )


def platform_identity_label(platform: str, identity: dict[str, Any]) -> str:
    if platform == "kuaishou":
        return str(identity.get("eid") or identity.get("user_id") or "")
    if platform == "wechat_channels":
        return str(identity.get("username") or "")
    return str(identity.get("sec_user_id") or identity.get("unique_id") or identity.get("user_id") or "")


def account_matches(account: dict[str, Any], args: argparse.Namespace) -> bool:
    if args.platform and normalize_platform(args.platform) != account.get("platform"):
        return False
    values = {
        str(account.get("homepage_url") or "").strip().lower(),
        str(account.get("display_name") or "").strip().lower(),
        str((account.get("profile_snapshot") or {}).get("display_name") or "").strip().lower(),
    }
    for value in (account.get("identity") or {}).values():
        values.add(str(value or "").strip().lower())
    if args.account:
        return args.account.strip().lower() in values
    if args.sec_user_id:
        return str((account.get("identity") or {}).get("sec_user_id") or "").strip().lower() == args.sec_user_id.strip().lower()
    if args.user_id:
        return str((account.get("identity") or {}).get("user_id") or "").strip().lower() == args.user_id.strip().lower()
    if args.unique_id:
        return str((account.get("identity") or {}).get("unique_id") or "").strip().lower() == args.unique_id.strip().lower()
    if args.username:
        return str((account.get("identity") or {}).get("username") or "").strip().lower() == args.username.strip().lower()
    if args.eid:
        return str((account.get("identity") or {}).get("eid") or "").strip().lower() == args.eid.strip().lower()
    return True


def direct_account_from_args(args: argparse.Namespace) -> dict[str, Any] | None:
    identity = {
        "sec_user_id": args.sec_user_id or "",
        "user_id": args.user_id or "",
        "unique_id": args.unique_id or "",
        "username": args.username or "",
        "eid": args.eid or "",
    }
    if not any(identity.values()):
        return None
    platform = normalize_platform(args.platform or ("wechat_channels" if args.username else "kuaishou" if args.eid else "douyin"))
    return {
        "platform": platform,
        "display_name": args.account or platform_identity_label(platform, identity) or platform_label(platform),
        "identity": {key: value for key, value in identity.items() if value},
        "enabled": True,
        "_direct": True,
    }


def account_retry_command(workspace: Path, account: dict[str, Any]) -> str:
    platform = str(account.get("platform") or "douyin")
    identity = dict(account.get("identity") or {})
    base = f"python3 douyin-vibe-track/scripts/run_daily.py --workspace {shell_quote(str(workspace))} --platform {shell_quote(platform)}"
    if platform == "wechat_channels" and identity.get("username"):
        return f"{base} --username {shell_quote(str(identity['username']))}"
    if platform == "kuaishou" and identity.get("eid"):
        return f"{base} --eid {shell_quote(str(identity['eid']))}"
    if identity.get("sec_user_id"):
        return f"{base} --sec-user-id {shell_quote(str(identity['sec_user_id']))}"
    if identity.get("unique_id"):
        return f"{base} --unique-id {shell_quote(str(identity['unique_id']))}"
    if identity.get("user_id"):
        return f"{base} --user-id {shell_quote(str(identity['user_id']))}"
    return f"{base} --account {shell_quote(account_label(account))}"


def filter_accounts(accounts: list[dict[str, Any]], args: argparse.Namespace, workspace: Path) -> tuple[list[dict[str, Any]], list[RunIssue]]:
    issues: list[RunIssue] = []
    selected = [account for account in accounts if account.get("enabled", True)]
    direct = direct_account_from_args(args)
    if direct:
        return [direct], issues
    if args.account or args.platform:
        matched = [account for account in selected if account_matches(account, args)]
        if not matched:
            issues.append(
                RunIssue(
                    phase="select_account",
                    platform=args.platform or "",
                    account=args.account or args.username or args.eid or args.sec_user_id or args.user_id or args.unique_id,
                    message="no enabled account matched the provided selector in accounts.json",
                    fix="Check accounts.json, then retry with --account using display_name, homepage_url, or the platform identity. If the user is not configured, first resolve and confirm the account profile before rerunning.",
                    command=f"python3 douyin-vibe-track/scripts/run_daily.py --workspace {shell_quote(str(workspace))} --account <display_name_or_identity>",
                )
            )
            return [], issues
        return matched, issues
    return selected, issues


def sanitize_filename(value: str) -> str:
    value = re.sub(r'[\\/:*?"<>|\s]+', "_", value).strip("._ ")
    return value[:120] or "video"


def account_data_path(accounts_dir: Path, account: dict[str, Any]) -> Path:
    platform = str(account.get("platform") or "douyin")
    identity = platform_identity_label(platform, account.get("identity") or {}) or account_label(account)
    return accounts_dir / f"{sanitize_filename(platform)}-{sanitize_filename(identity)}.json"


def candidate_from_dict(data: dict[str, Any]) -> Candidate:
    metrics = data.get("metrics") or {}
    video = data.get("video") or {}
    if "video_urls" in data and "download_urls" not in video:
        video["download_urls"] = data.get("video_urls") or []
    return Candidate(
        platform=str(data.get("platform") or "douyin"),
        post_id=str(data.get("post_id") or data.get("aweme_id") or ""),
        author_name=str(data.get("author_name") or ""),
        author_id=str(data.get("author_id") or data.get("sec_user_id") or ""),
        desc=str(data.get("desc") or ""),
        share_url=str(data.get("share_url") or ""),
        create_time=int(data.get("create_time") or 0),
        metrics={
            "like_count": int(metrics.get("like_count") or data.get("digg_count") or 0),
            "collect_count": int(metrics.get("collect_count") or data.get("collect_count") or 0),
            "share_count": int(metrics.get("share_count") or data.get("share_count") or 0),
            "comment_count": int(metrics.get("comment_count") or data.get("comment_count") or 0),
        },
        video={
            "download_urls": [str(url) for url in video.get("download_urls") or data.get("video_urls") or []],
            "decrypt_key": str(video.get("decrypt_key") or ""),
            "needs_decrypt": bool(video.get("needs_decrypt")),
        },
        video_path=str(data.get("video_path") or ""),
    )


def read_completed_account_candidates(path: Path) -> list[Candidate] | None:
    try:
        payload = load_json(path, {})
    except Exception:
        return None
    if payload.get("status") != "completed":
        return None
    return [candidate_from_dict(item) for item in payload.get("candidates") or []]


def write_account_candidates(path: Path, account: dict[str, Any], candidates: list[Candidate], generated_at: datetime) -> None:
    write_json(
        path,
        {
            "status": "completed",
            "generated_at": generated_at.isoformat(),
            "account": {
                "platform": account.get("platform") or "",
                "display_name": account.get("display_name") or "",
                "homepage_url": account.get("homepage_url") or "",
                "identity": account.get("identity") or {},
                "profile_snapshot": account.get("profile_snapshot") or {},
            },
            "candidates": [asdict(candidate) for candidate in candidates],
        },
    )


def download_file(urls: list[str], destination: Path, timeout: int | tuple[int, int], max_retries: int) -> Path:
    destination.parent.mkdir(parents=True, exist_ok=True)
    last_error: Exception | None = None
    session = get_session()
    for url in urls:
        for attempt in range(1, max_retries + 1):
            try:
                with session.get(url, stream=True, timeout=timeout, headers={"User-Agent": "Mozilla/5.0"}) as response:
                    response.raise_for_status()
                    tmp = destination.with_suffix(destination.suffix + ".tmp")
                    with tmp.open("wb") as fh:
                        for chunk in response.iter_content(chunk_size=1024 * 1024):
                            if chunk:
                                fh.write(chunk)
                    tmp.replace(destination)
                    return destination
            except Exception as exc:  # noqa: BLE001
                last_error = exc
                if attempt < max_retries and should_retry_exception(exc):
                    time.sleep(min(2 ** (attempt - 1), 8))
                    continue
                break
    raise RuntimeError(f"download failed after trying {len(urls)} URLs: {last_error}") from last_error


def decrypt_wechat_channels_video(source: Path, destination: Path, decode_key: str) -> Path:
    node = shutil.which("node")
    if not node:
        raise RuntimeError("wechat_channels decryption requires Node.js in PATH")
    helper = SCRIPT_DIR / "wechat_channels_keystream.js"
    result = subprocess.run(
        [node, str(helper), decode_key, "hex"],
        check=True,
        capture_output=True,
        text=True,
    )
    keystream = bytes.fromhex(result.stdout.strip())
    destination.parent.mkdir(parents=True, exist_ok=True)
    with source.open("rb") as src, destination.open("wb") as dst:
        prefix = src.read(len(keystream))
        decrypted = bytes(byte ^ keystream[index] for index, byte in enumerate(prefix))
        dst.write(decrypted)
        shutil.copyfileobj(src, dst)
    return destination


def normalize_video_for_ppt(source: Path, destination: Path) -> Path:
    ffmpeg = shutil.which("ffmpeg")
    if not ffmpeg:
        if source != destination:
            source.replace(destination)
        return destination
    tmp = destination.with_name(f"{destination.stem}.ppt-tmp.mp4")
    command = [
        ffmpeg,
        "-y",
        "-i",
        str(source),
        "-map",
        "0:v:0",
        "-map",
        "0:a?",
        "-c:v",
        "libx264",
        "-profile:v",
        "main",
        "-pix_fmt",
        "yuv420p",
        "-preset",
        "veryfast",
        "-crf",
        "23",
        "-c:a",
        "aac",
        "-b:a",
        "128k",
        "-movflags",
        "+faststart",
        str(tmp),
    ]
    subprocess.run(command, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.PIPE)
    tmp.replace(destination)
    if source != destination and source.exists():
        source.unlink()
    return destination


def extract_poster_frame(video_path: Path) -> Path | None:
    ffmpeg = shutil.which("ffmpeg")
    if not ffmpeg:
        return None
    poster_path = video_path.with_suffix(".jpg")
    commands = [
        [ffmpeg, "-y", "-i", str(video_path), "-frames:v", "1", "-q:v", "2", str(poster_path)],
        [ffmpeg, "-y", "-i", str(video_path), "-vf", "thumbnail", "-frames:v", "1", "-q:v", "2", str(poster_path)],
    ]
    for command in commands:
        poster_path.unlink(missing_ok=True)
        try:
            subprocess.run(command, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.PIPE)
        except subprocess.CalledProcessError:
            continue
        if poster_path.exists():
            return poster_path
    return None


def format_time(epoch_seconds: int, tz: ZoneInfo) -> str:
    if not epoch_seconds:
        return ""
    return datetime.fromtimestamp(epoch_seconds, tz).strftime("%Y-%m-%d %H:%M:%S")


def format_filename_time(epoch_seconds: int, tz: ZoneInfo) -> str:
    if not epoch_seconds:
        return "unknown-time"
    return datetime.fromtimestamp(epoch_seconds, tz).strftime("%Y-%m-%d_%H-%M-%S")


def make_video_destination(videos_dir: Path, item: Candidate, tz: ZoneInfo) -> Path:
    author = sanitize_filename(item.author_name or item.author_id or item.post_id)
    created_at = format_filename_time(item.create_time, tz)
    platform = sanitize_filename(item.platform)
    destination = videos_dir / f"{platform}-{author}-{created_at}.mp4"
    if not destination.exists():
        return destination
    for index in range(2, 1000):
        candidate = videos_dir / f"{platform}-{author}-{created_at}-{index}.mp4"
        if not candidate.exists():
            return candidate
    raise RuntimeError(f"could not find available filename for {platform}-{author}-{created_at}")


def load_pptx_dependencies() -> tuple[Any, Any, Any]:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as exc:  # pragma: no cover
        raise RuntimeError("Missing dependency python-pptx. Install with: python3 -m pip install -r requirements.txt") from exc
    return Presentation, Inches, Pt


def build_pptx(path: Path, candidates: list[Candidate], tz: ZoneInfo) -> None:
    Presentation, Inches, Pt = load_pptx_dependencies()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    poster_frames: list[Path] = []

    if not candidates:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.8), Inches(1.5))
        text = box.text_frame
        text.text = "今日无点赞超过阈值的新作品"
        text.paragraphs[0].font.size = Pt(34)
        prs.save(path)
        return

    for item in candidates:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        video_path = Path(item.video_path)
        if video_path.exists():
            poster_frame = extract_poster_frame(video_path)
            if poster_frame:
                poster_frames.append(poster_frame)
            slide.shapes.add_movie(
                str(video_path),
                Inches(0.35),
                Inches(0.35),
                Inches(6.2),
                Inches(6.8),
                poster_frame_image=str(poster_frame) if poster_frame else None,
                mime_type="video/mp4",
            )

        title_box = slide.shapes.add_textbox(Inches(6.9), Inches(0.55), Inches(5.8), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = item.author_name or item.author_id or item.post_id
        title_frame.paragraphs[0].font.size = Pt(26)
        title_frame.paragraphs[0].font.bold = True

        details = [
            ("平台", platform_label(item.platform)),
            ("博主名称", item.author_name or item.author_id or item.post_id),
            ("点赞数", item.like_count),
            ("收藏数", item.collect_count),
            ("分享数", item.share_count),
            ("发布时间", format_time(item.create_time, tz)),
        ]
        body = slide.shapes.add_textbox(Inches(6.9), Inches(1.55), Inches(5.8), Inches(4.9)).text_frame
        body.word_wrap = True
        body.text = ""
        for label, value in details:
            paragraph = body.add_paragraph()
            paragraph.text = f"{label}: {value}"
            paragraph.font.size = Pt(18)

    prs.save(path)
    for poster_frame in poster_frames:
        poster_frame.unlink(missing_ok=True)


def print_issues(issues: list[RunIssue]) -> None:
    for issue in issues:
        prefix_parts = [issue.phase]
        if issue.platform:
            prefix_parts.append(f"platform={issue.platform}")
        if issue.account:
            prefix_parts.append(f"account={issue.account}")
        if issue.post_id:
            prefix_parts.append(f"post_id={issue.post_id}")
        print(f"[ERROR] {' '.join(prefix_parts)}: {issue.message}", file=sys.stderr)
        print(f"[FIX] {issue.fix}", file=sys.stderr)
        if issue.command:
            print(f"[RETRY] {issue.command}", file=sys.stderr)


def extract_kuaishou_video_urls(item: dict[str, Any]) -> list[str]:
    urls: list[str] = []
    photo = (item.get("photo") or {}) if item.get("photo") else item
    cover = photo.get("photoUrl") or photo.get("playUrl") or ""
    for value in [cover]:
        if isinstance(value, str) and value and value not in urls:
            urls.append(value)
    for key in ("mainMvUrls", "photoH265Urls", "photoH264Urls"):
        for entry in photo.get(key) or []:
            url = entry.get("url")
            if url and url not in urls:
                urls.append(str(url))
    manifest = photo.get("manifest") or {}
    for adaptation_set in manifest.get("adaptationSet") or []:
        for representation in adaptation_set.get("representation") or []:
            url = representation.get("url")
            if url and url not in urls:
                urls.append(str(url))
    return urls


def extract_wechat_video_urls(item: dict[str, Any]) -> tuple[list[str], str]:
    object_desc = item.get("object_desc") or item.get("objectDesc") or {}
    media_items = object_desc.get("media") or object_desc.get("media_list") or []
    urls: list[str] = []
    decode_key = ""
    for media in media_items:
        url = str(media.get("url") or media.get("play_url") or "")
        token = str(media.get("url_token") or media.get("token") or "")
        decode_key = decode_key or str(media.get("decode_key") or "")
        if url:
            full_url = f"{url}{token}" if token and token not in url else url
            if full_url not in urls:
                urls.append(full_url)
    if not decode_key:
        decode_key = str(object_desc.get("decode_key") or item.get("decode_key") or "")
    return urls, decode_key


def to_candidate(platform: str, item: dict[str, Any]) -> Candidate | None:
    platform = normalize_platform(platform)
    if platform in {"douyin", "tiktok"}:
        author = item.get("author") or {}
        stats = item.get("statistics") or {}
        video = item.get("video") or {}
        urls: list[str] = []
        for key in ("download_no_watermark_addr", "play_addr", "download_addr"):
            address = video.get(key) or {}
            for url in address.get("url_list") or []:
                if url and url not in urls:
                    urls.append(str(url))
        post_id = str(item.get("aweme_id") or item.get("id") or "")
        if not post_id:
            return None
        return Candidate(
            platform=platform,
            post_id=post_id,
            author_name=str(author.get("nickname") or ""),
            author_id=str(author.get("sec_uid") or author.get("unique_id") or ""),
            desc=str(item.get("desc") or ""),
            share_url=str(item.get("share_url") or ""),
            create_time=int(item.get("create_time") or 0),
            metrics={
                "like_count": int(stats.get("digg_count") or 0),
                "collect_count": int(stats.get("collect_count") or 0),
                "share_count": int(stats.get("share_count") or 0),
                "comment_count": int(stats.get("comment_count") or 0),
            },
            video={"download_urls": urls, "decrypt_key": "", "needs_decrypt": False},
        )

    if platform == "kuaishou":
        photo = (item.get("photo") or {}) if item.get("photo") else item
        author = photo.get("author") or item.get("author") or {}
        post_id = str(photo.get("photoId") or photo.get("id") or "")
        if not post_id:
            return None
        urls = extract_kuaishou_video_urls(item)
        return Candidate(
            platform=platform,
            post_id=post_id,
            author_name=str(author.get("name") or author.get("user_name") or ""),
            author_id=str(author.get("id") or author.get("eid") or ""),
            desc=str(photo.get("caption") or photo.get("title") or ""),
            share_url=str(photo.get("shareUrl") or photo.get("photoUrl") or ""),
            create_time=int(photo.get("timestamp") or photo.get("createTime") or 0),
            metrics={
                "like_count": summarize_counts(photo.get("realLikeCount") or photo.get("likeCount")),
                "collect_count": summarize_counts(photo.get("viewCount")),
                "share_count": summarize_counts(photo.get("shareCount")),
                "comment_count": summarize_counts(photo.get("commentCount")),
            },
            video={"download_urls": urls, "decrypt_key": "", "needs_decrypt": False},
        )

    object_desc = item.get("object_desc") or item.get("objectDesc") or {}
    contact = item.get("contact") or {}
    ext_info = item.get("ext_info") or object_desc.get("ext_info") or (contact.get("ext_info") or {})
    username = str(item.get("username") or item.get("finder_username") or contact.get("username") or "")
    post_id = str(item.get("object_id") or item.get("id") or object_desc.get("object_id") or ext_info.get("feed_id") or "")
    if not post_id:
        return None
    urls, decode_key = extract_wechat_video_urls(item)
    create_time = (
        item.get("create_time")
        or item.get("createtime")
        or object_desc.get("create_time")
        or ext_info.get("publish_time")
        or ext_info.get("create_time")
        or 0
    )
    return Candidate(
        platform=platform,
        post_id=post_id,
        author_name=str(item.get("nickname") or contact.get("nickname") or item.get("name") or username),
        author_id=username,
        desc=str(object_desc.get("description") or object_desc.get("desc") or ""),
        share_url=str(item.get("share_url") or object_desc.get("share_url") or ""),
        create_time=int(create_time or 0),
        metrics={
            "like_count": summarize_counts(item.get("like_count") or ext_info.get("like_num") or item.get("like_num") or object_desc.get("like_num")),
            "collect_count": summarize_counts(item.get("fav_count") or ext_info.get("fav_num") or object_desc.get("fav_num")),
            "share_count": summarize_counts(item.get("forward_count") or ext_info.get("forward_num") or ext_info.get("share_num") or object_desc.get("forward_num")),
            "comment_count": summarize_counts(item.get("comment_count") or ext_info.get("comment_num") or object_desc.get("comment_num")),
        },
        video={"download_urls": urls, "decrypt_key": decode_key, "needs_decrypt": bool(decode_key)},
    )


def iter_recent_posts(
    account: dict[str, Any],
    config: dict[str, Any],
    since_epoch: int,
    workspace: Path,
) -> tuple[list[dict[str, Any]], list[RunIssue], int]:
    platform = str(account.get("platform") or "douyin")
    identity = account.get("identity") or {}
    posts: list[dict[str, Any]] = []
    errors: list[RunIssue] = []
    cursor: str | int | None = 0 if platform in {"douyin", "tiktok"} else ""
    last_cursor = None

    pages_fetched = 0
    for _ in range(int(config.get("max_pages_per_account") or 3)):
        try:
            payload = fetch_posts_page(workspace, platform, identity, cursor)
            payload["_cursor"] = cursor
            page_posts, next_cursor, has_more = unwrap_posts_page(platform, payload)
            pages_fetched += 1
        except Exception as exc:  # noqa: BLE001
            errors.append(
                RunIssue(
                    phase="fetch_posts",
                    platform=platform,
                    account=account_label(account),
                    message=str(exc),
                    fix="Retry this account only. If it still fails, verify the stored platform identity, check the TikHub API key/quota, and increase request_timeout_seconds or max_retries in config.json for transient failures.",
                    command=account_retry_command(workspace, account),
                )
            )
            break

        if not page_posts:
            break
        posts.extend(page_posts)
        if platform == "kuaishou":
            break

        candidate_times = [int((to_candidate(platform, item) or Candidate(platform, "", "", "", "", "", 0, {}, {})).create_time or 0) for item in page_posts]
        oldest = min([value for value in candidate_times if value], default=0)
        if oldest and oldest < since_epoch:
            break
        if not has_more or next_cursor in (None, "", last_cursor, cursor):
            break
        last_cursor = cursor
        cursor = next_cursor
    return posts, errors, pages_fetched


def collect_account_candidates(
    account: dict[str, Any],
    config: dict[str, Any],
    since_epoch: int,
    workspace: Path,
    reported_ids: set[str],
) -> AccountFetchResult:
    started = time.perf_counter()
    posts, account_errors, pages_fetched = iter_recent_posts(account, config, since_epoch, workspace)
    if account_errors:
        return AccountFetchResult(account, [], account_errors, time.perf_counter() - started, len(posts), pages_fetched)

    account_candidates: list[Candidate] = []
    account_issues: list[RunIssue] = []
    seen_account: set[str] = set()
    threshold = platform_threshold(config, str(account.get("platform") or "douyin"))
    platform = str(account.get("platform") or "douyin")

    for post in posts:
        candidate = to_candidate(platform, post)
        if not candidate or not candidate.post_id:
            continue
        candidate_key = post_key(candidate.platform, candidate.post_id)
        if candidate_key in reported_ids or candidate_key in seen_account:
            continue
        if candidate.create_time and candidate.create_time < since_epoch:
            continue
        if candidate.like_count < threshold:
            continue
        if not candidate.video_urls:
            account_issues.append(
                RunIssue(
                    phase="collect_video_url",
                    platform=candidate.platform,
                    account=candidate.author_name or account_label(account),
                    post_id=candidate.post_id,
                    message="no downloadable video URL found in TikHub response",
                    fix="Retry this account only. If the response still has no video URL, inspect the raw TikHub response for this post and switch to another available video source if needed.",
                    command=account_retry_command(workspace, account),
                )
            )
            continue
        if candidate.needs_decrypt and not candidate.decrypt_key:
            account_issues.append(
                RunIssue(
                    phase="collect_decrypt_key",
                    platform=candidate.platform,
                    account=candidate.author_name or account_label(account),
                    post_id=candidate.post_id,
                    message="wechat_channels video requires a decode_key but none was present",
                    fix="Retry this account only. If it still fails, inspect the raw TikHub object_desc payload and verify the endpoint still returns decode_key for downloadable videos.",
                    command=account_retry_command(workspace, account),
                )
            )
            continue
        seen_account.add(candidate_key)
        account_candidates.append(candidate)

    if account_issues:
        return AccountFetchResult(account, [], account_issues, time.perf_counter() - started, len(posts), pages_fetched)
    account_candidates.sort(key=lambda item: (item.like_count, item.create_time), reverse=True)
    return AccountFetchResult(account, account_candidates, [], time.perf_counter() - started, len(posts), pages_fetched)


def process_candidate_download(
    item: Candidate,
    config: dict[str, Any],
    workspace: Path,
    videos_dir: Path,
    tz: ZoneInfo,
    timeout: tuple[int, int],
    max_retries: int,
) -> DownloadResult:
    started = time.perf_counter()
    destination = make_video_destination(videos_dir, item, tz)
    raw_destination = destination.with_name(f"{destination.stem}.raw.mp4")
    normalized_source = raw_destination
    try:
        download_file(item.video_urls, raw_destination, timeout, max_retries)
        if item.platform == "wechat_channels" and item.needs_decrypt:
            decrypted_destination = destination.with_name(f"{destination.stem}.decrypted.mp4")
            decrypt_wechat_channels_video(raw_destination, decrypted_destination, item.decrypt_key)
            keep_encrypted = bool(((config.get("wechat_channels") or {}).get("keep_encrypted_copy")) or False)
            if not keep_encrypted:
                raw_destination.unlink(missing_ok=True)
            normalized_source = decrypted_destination
        normalize_video_for_ppt(normalized_source, destination)
        item.video_path = str(destination)
        return DownloadResult(item, None, time.perf_counter() - started)
    except Exception as exc:  # noqa: BLE001
        return DownloadResult(
            item,
            RunIssue(
                phase="download_video" if item.platform != "wechat_channels" or not item.needs_decrypt else "decrypt_video",
                platform=item.platform,
                account=item.author_name or item.author_id or platform_label(item.platform),
                post_id=item.post_id,
                message=str(exc),
                fix="Retry this account only. If it still fails, verify the video URL and, for 视频号, confirm Node.js and the decode_key flow are available before rerunning.",
                command=(
                    f"python3 douyin-vibe-track/scripts/run_daily.py --workspace {shell_quote(str(workspace))} --platform {shell_quote(item.platform)} --username {shell_quote(item.author_id)}"
                    if item.platform == "wechat_channels" and item.author_id
                    else f"python3 douyin-vibe-track/scripts/run_daily.py --workspace {shell_quote(str(workspace))} --platform {shell_quote(item.platform)} --account {shell_quote(item.author_name or item.author_id)}"
                ),
            ),
            time.perf_counter() - started,
        )


def platform_threshold(config: dict[str, Any], platform: str) -> int:
    thresholds = config.get("thresholds") or {}
    platform_block = thresholds.get(platform) or {}
    return int(platform_block.get("like_count") or DEFAULT_THRESHOLDS.get(platform, {"like_count": 10000})["like_count"])


def post_key(platform: str, post_id: str) -> str:
    return f"{platform}:{post_id}"


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Run one daily multi-platform monitoring pass.")
    parser.add_argument("--workspace", default=str(Path.home() / WORKSPACE_DIRNAME))
    parser.add_argument("--platform", choices=("douyin", "tiktok", "wechat_channels", "kuaishou"), default="")
    parser.add_argument("--account", default="", help="Run only one configured account, matched by display_name, homepage_url, or identity.")
    parser.add_argument("--sec-user-id", default="", help="Run one douyin/tiktok user directly by sec_user_id.")
    parser.add_argument("--user-id", default="", help="Run one user directly by user_id, or match a configured account by user_id.")
    parser.add_argument("--unique-id", default="", help="Run one douyin/tiktok user directly by unique_id.")
    parser.add_argument("--username", default="", help="Run one wechat_channels user directly by username.")
    parser.add_argument("--eid", default="", help="Run one kuaishou user directly by eid.")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    workspace = Path(args.workspace).expanduser()
    config_path = workspace / "config.json"
    accounts_path = workspace / "accounts.json"
    state_path = workspace / "state.json"

    try:
        config = load_json(config_path, {})
        accounts_data = load_json(accounts_path, {"accounts": []})
        state = load_json(state_path, {"last_run_at": None, "reported_aweme_ids": []})
    except Exception as exc:  # noqa: BLE001
        print_issues(
            [
                RunIssue(
                    phase="load_workspace",
                    message=str(exc),
                    fix="Check config.json, accounts.json, and state.json for invalid JSON or unreadable files. Repair the file named in the error, then rerun the command.",
                )
            ]
        )
        return 2

    config, config_changed = migrate_config(config)
    accounts_data, accounts_changed = migrate_accounts(accounts_data)
    state, state_changed = migrate_state(state)
    if config_changed:
        write_json(config_path, config)
    if accounts_changed:
        write_json(accounts_path, accounts_data)
    if state_changed:
        write_json(state_path, state)

    api_key = config.get("tikhub_api_key") or ""
    if not api_key:
        print_issues(
            [
                RunIssue(
                    phase="config",
                    message=f"missing tikhub_api_key in {config_path}",
                    fix="Configure config.json with a valid TikHub API key before running. Do not prompt inside the script; the Agent should ask the user for the key or use an existing configured key.",
                )
            ]
        )
        return 2

    platform_enabled = {
        "douyin": True,
        "tiktok": True,
        "wechat_channels": bool(((config.get("platforms") or {}).get("wechat_channels") or {}).get("enabled", True)),
        "kuaishou": bool(((config.get("platforms") or {}).get("kuaishou") or {}).get("enabled", True)),
    }
    configured_enabled_accounts = [
        account
        for account in accounts_data.get("accounts", [])
        if account.get("enabled", True) and platform_enabled.get(str(account.get("platform") or "douyin"), True)
    ]
    enabled_accounts, selection_issues = filter_accounts(configured_enabled_accounts, args, workspace)
    if selection_issues:
        print_issues(selection_issues)
        print(json.dumps({"status": "failed", "errors": [issue.to_dict() for issue in selection_issues]}, ensure_ascii=False, indent=2))
        return 2
    if not enabled_accounts:
        print_issues(
            [
                RunIssue(
                    phase="select_account",
                    message=f"no enabled accounts in {accounts_path}",
                    fix="Add or enable at least one account in accounts.json. For new platforms, first resolve the account identity and confirm the profile before rerunning the monitor.",
                )
            ]
        )
        return 2

    try:
        load_pptx_dependencies()
    except RuntimeError as exc:
        print_issues(
            [
                RunIssue(
                    phase="dependencies",
                    message=str(exc),
                    fix="Install dependencies with: python3 -m pip install -r douyin-vibe-track/requirements.txt",
                )
            ]
        )
        return 2

    tz = ZoneInfo(config.get("timezone") or "Asia/Shanghai")
    now = datetime.now(tz)
    report_date = now.strftime("%Y-%m-%d")
    lookback_hours = int(config.get("lookback_hours") or 72)
    since_epoch = int((now - timedelta(hours=lookback_hours)).timestamp())
    max_retries = int(config.get("max_retries") or 3)
    connect_timeout = int(config.get("connect_timeout_seconds") or 5)
    read_timeout = int(config.get("read_timeout_seconds") or config.get("request_timeout_seconds") or 30)
    download_timeout = int(config.get("download_timeout_seconds") or max(read_timeout, 120))
    fetch_concurrency = max(1, int(config.get("fetch_concurrency") or 6))
    download_concurrency = max(1, int(config.get("download_concurrency") or 3))
    slow_account_threshold = float(config.get("slow_account_threshold_seconds") or 45)
    reported_ids = {str(value) for value in state.get("reported_post_keys", [])}
    metrics = RuntimeMetrics()
    overall_started = time.perf_counter()

    reports_root = resolve_reports_root(workspace, config)
    report_dir = reports_root / report_date
    accounts_dir = report_dir / "accounts"
    videos_dir = report_dir / "videos"
    accounts_dir.mkdir(parents=True, exist_ok=True)
    videos_dir.mkdir(parents=True, exist_ok=True)

    is_targeted_run = bool(
        args.account
        or args.sec_user_id
        or args.user_id
        or args.unique_id
        or args.username
        or args.eid
    )
    accounts_for_report = enabled_accounts if is_targeted_run or any(account.get("_direct") for account in enabled_accounts) else configured_enabled_accounts
    errors: list[RunIssue] = []
    pending_accounts = []
    for account in enabled_accounts:
        data_path = account_data_path(accounts_dir, account)
        if not is_targeted_run and read_completed_account_candidates(data_path) is not None:
            metrics.account_timings.append(
                {"account": account_label(account), "platform": str(account.get("platform") or "douyin"), "duration_seconds": 0.0, "status": "cached"}
            )
            continue
        pending_accounts.append(account)

    fetch_started = time.perf_counter()
    with concurrent.futures.ThreadPoolExecutor(max_workers=min(fetch_concurrency, max(1, len(pending_accounts)))) as executor:
        future_map = {
            executor.submit(collect_account_candidates, account, config, since_epoch, workspace, reported_ids): account
            for account in pending_accounts
        }
        for future in concurrent.futures.as_completed(future_map):
            result = future.result()
            metrics.fetched_accounts += 1
            status = "completed"
            if result.errors:
                metrics.failed_accounts += 1
                errors.extend(result.errors)
                status = "failed"
            else:
                write_account_candidates(account_data_path(accounts_dir, result.account), result.account, result.candidates, now)
            timing = {
                "account": account_label(result.account),
                "platform": str(result.account.get("platform") or "douyin"),
                "duration_seconds": round(result.duration_seconds, 3),
                "posts_seen": result.posts_seen,
                "pages_fetched": result.pages_fetched,
                "status": status,
            }
            metrics.account_timings.append(timing)
            if result.duration_seconds >= slow_account_threshold:
                metrics.slow_accounts.append(timing)
    metrics.fetch_stage_seconds = round(time.perf_counter() - fetch_started, 3)

    candidates: list[Candidate] = []
    missing_accounts: list[dict[str, Any]] = []
    for account in accounts_for_report:
        account_candidates = read_completed_account_candidates(account_data_path(accounts_dir, account))
        if account_candidates is None:
            missing_accounts.append(account)
            continue
        candidates.extend(account_candidates)

    existing_error_accounts = {issue.account for issue in errors if issue.account}
    for account in missing_accounts:
        if account_label(account) in existing_error_accounts:
            continue
        errors.append(
            RunIssue(
                phase="account_data",
                platform=str(account.get("platform") or ""),
                account=account_label(account),
                message=f"missing completed account data file: {account_data_path(accounts_dir, account)}",
                fix="Fetch this account only. The daily PPT will be generated only after every enabled account has a completed JSON data file for today.",
                command=account_retry_command(workspace, account),
            )
        )

    if errors:
        issue_payload = {
            "status": "needs_repair",
            "generated_at": now.isoformat(),
            "issues": [issue.to_dict() for issue in errors],
            "recovery_commands": sorted({issue.command for issue in errors if issue.command}),
            "metrics": asdict(metrics),
        }
        write_json(report_dir / "issues.json", issue_payload)
        print_issues(errors)
        print(
            json.dumps(
                {
                    "status": "needs_repair",
                    "report_dir": str(report_dir),
                    "account_data_dir": str(accounts_dir),
                    "pptx": str(report_dir / "summary.pptx"),
                    "videos": 0,
                    "errors": len(errors),
                    "issues": [issue.to_dict() for issue in errors],
                    "recovery_commands": issue_payload["recovery_commands"],
                    "metrics": asdict(metrics),
                },
                ensure_ascii=False,
                indent=2,
            )
        )
        return 1

    candidates.sort(key=lambda item: (item.like_count, item.create_time), reverse=True)
    metrics.download_candidates = len(candidates)

    download_queue: list[Candidate] = []
    for item in candidates:
        video_stem = f"{sanitize_filename(item.platform)}-{sanitize_filename(item.author_name or item.author_id or item.post_id)}-{format_filename_time(item.create_time, tz)}"
        existing_video = next(
            (
                path
                for path in videos_dir.iterdir()
                if path.suffix == ".mp4"
                and not path.name.endswith(".raw.mp4")
                and ".ppt-tmp." not in path.name
                and (path.stem == video_stem or path.stem.startswith(f"{video_stem}-"))
            ),
            None,
        )
        if existing_video and not is_targeted_run:
            item.video_path = str(existing_video)
            continue
        download_queue.append(item)

    download_started = time.perf_counter()
    download_timeout_tuple = (connect_timeout, download_timeout)
    with concurrent.futures.ThreadPoolExecutor(max_workers=min(download_concurrency, max(1, len(download_queue)))) as executor:
        future_map = {
            executor.submit(
                process_candidate_download,
                item,
                config,
                workspace,
                videos_dir,
                tz,
                download_timeout_tuple,
                max_retries,
            ): item
            for item in download_queue
        }
        for future in concurrent.futures.as_completed(future_map):
            result = future.result()
            if result.error:
                errors.append(result.error)
            elif result.item.video_path:
                metrics.downloaded_videos += 1
    metrics.download_stage_seconds = round(time.perf_counter() - download_started, 3)
    metrics.downloaded_videos += sum(1 for item in candidates if item.video_path) - metrics.downloaded_videos

    written_candidates = [item for item in candidates if item.video_path]
    if errors:
        issue_payload = {
            "status": "needs_repair",
            "generated_at": now.isoformat(),
            "issues": [issue.to_dict() for issue in errors],
            "recovery_commands": sorted({issue.command for issue in errors if issue.command}),
            "metrics": asdict(metrics),
        }
        write_json(report_dir / "issues.json", issue_payload)
        print_issues(errors)
        print(
            json.dumps(
                {
                    "status": "needs_repair",
                    "report_dir": str(report_dir),
                    "account_data_dir": str(accounts_dir),
                    "pptx": str(report_dir / "summary.pptx"),
                    "videos": len(written_candidates),
                    "errors": len(errors),
                    "issues": [issue.to_dict() for issue in errors],
                    "recovery_commands": issue_payload["recovery_commands"],
                    "metrics": asdict(metrics),
                },
                ensure_ascii=False,
                indent=2,
            )
        )
        return 1

    pptx_built = False
    ppt_started = time.perf_counter()
    try:
        build_pptx(report_dir / "summary.pptx", written_candidates, tz)
        pptx_built = True
    except Exception as exc:  # noqa: BLE001
        errors.append(
            RunIssue(
                phase="build_pptx",
                message=str(exc),
                fix="The videos may already be downloaded. Check the listed video paths, install or repair python-pptx/ffmpeg if needed, then rerun the same command after fixing the PPT generation error.",
                command=f"python3 douyin-vibe-track/scripts/run_daily.py --workspace {shell_quote(str(workspace))}",
            )
        )
    metrics.ppt_stage_seconds = round(time.perf_counter() - ppt_started, 3)

    state["last_run_at"] = now.isoformat()
    state["reported_post_keys"] = (
        sorted(reported_ids | {post_key(item.platform, item.post_id) for item in written_candidates}) if pptx_built else sorted(reported_ids)
    )
    write_json(state_path, state)

    if errors:
        issue_payload = {
            "status": "needs_repair",
            "generated_at": now.isoformat(),
            "issues": [issue.to_dict() for issue in errors],
            "recovery_commands": sorted({issue.command for issue in errors if issue.command}),
            "metrics": asdict(metrics),
        }
        write_json(report_dir / "issues.json", issue_payload)
        print_issues(errors)
    recovery_commands = sorted({issue.command for issue in errors if issue.command})
    metrics.total_seconds = round(time.perf_counter() - overall_started, 3)
    metrics.slow_accounts = sorted(metrics.slow_accounts, key=lambda item: item["duration_seconds"], reverse=True)[:5]
    print(
        json.dumps(
            {
                "status": "needs_repair" if errors else "completed",
                "report_dir": str(report_dir),
                "pptx": str(report_dir / "summary.pptx"),
                "videos": len(written_candidates),
                "errors": len(errors),
                "issues": [issue.to_dict() for issue in errors],
                "recovery_commands": recovery_commands,
                "metrics": asdict(metrics),
            },
            ensure_ascii=False,
            indent=2,
        )
    )
    return 0 if not errors else 1


if __name__ == "__main__":
    raise SystemExit(main())
